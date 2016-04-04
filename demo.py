from pptx import Presentation
from pptx.util import Length, Pt
from canvas import Canvas

from pprint import pprint

if __name__ == '__main__':

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    canvas = Canvas(slide, prs.slide_height, prs.slide_width)
    ctx = canvas.getContext('2d')

    ctx.fillStyle = '#C0C0C0'
    pprint(ctx.fillStyle)

    ctx.lineWidth = Pt(4)
    pprint(dict((k, ctx.__dict__['_' + k]) for k in ctx._keys if '_' + k in ctx.__dict__))



#    pprint(Length(ctx.lineWidth).pt())

    ctx.fillRect(Pt(100), Pt(100), Pt(150), Pt(250))

    prs.save('output.pptx')
