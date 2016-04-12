from pptx import Presentation
from pptx.util import Length, Pt
from os import sys, path
from math import atan2, ceil, cos, pi, pow, sin, sqrt

DATA = [{
    'name': 'test',
    'color': 'blue',
    'data': [
        0.7,
        1.7,
        1.5,
        2.3
    ]
}]

class Radar(object):

    def __init__(self, ctx, data, scale=0.8):
        self.ctx = ctx
        self.data = data
        self.scale = scale

    def sin(self, degrees):
        return sin(pi * degrees / 180)

    def cos(self, degrees):
        return cos(pi * degrees / 180)

    def quadrant(self, x1, y1, x2, y2):
        x = x2 - x1
        y = y2 - y1

        if x >= 0 and y >= 0:
            return 1
        if x <= 0 and y >= 0:
            return 2
        if x <= 0 and y <= 0:
            return 3
        if x >= 0 and y <= 0:
            return 4
        print x
        print y
        raise Exception

    def getVertices(self, x, y):
        radius = (x if (x < y) else y) * self.scale
        vertices = []

        for i in range(len(self.data)):
            degrees = i * (360 / len(self.data))
            new_x = self.sin(degrees) * radius + x
            new_y = self.cos(degrees) * radius + y

            vertices.append([ new_x, (y * 2) - new_y ])

        return vertices

    def draw(self, x, y):
        self.drawAxis(x, y)
        self.drawData(x, y)

    def drawAxis(self, x, y):
        vertices = self.getVertices(x, y)

        for i in range(len(self.data)):
            self.ctx.beginPath()
            self.ctx.moveTo(Pt(x), Pt(y))
            self.ctx.lineTo(Pt(vertices[i][0]), Pt(vertices[i][1]))
            self.ctx.strokeStyle = '#666'
            self.ctx.lineWidth = Pt(2)
            self.ctx.stroke()
            self.ctx.closePath()

            angle = atan2(vertices[i][1] - y, vertices[i][0] - x) * 180 / pi + 180;

            max_val = int(ceil(max(self.data)))
            if i is 0:
                formatter = '{0} = {0}x the average'
            else:
                formatter = '{0}'

            self.drawStub(formatter.format(max_val), vertices[i][0], vertices[i][1], angle)

            for j in range(max_val):
                if j is 0:
                  continue

                midway = self.getPoint(x, y, vertices[i][0], vertices[i][1], angle, (100 / max_val) * j)
                if i is 0 and j is 1:
                    text = '{0} = average score'
                else:
                    text = '{0}'
                self.drawStub(text.format(j), midway['x'], midway['y'], angle)

    def getPoint(self, x1, y1, x2, y2, angle, percent):
        distance = sqrt(pow(x2 - x1, 2) + pow(y2 - y1, 2))

        return {
            'x': x1 - self.cos(angle) * (distance * (percent / 100.0)),
            'y': y1 - self.sin(angle) * (distance * (percent / 100.0))
        }

    def drawStub(self, text, x, y, angle):
        stub_x1 = x + (5 * self.cos(angle + 90))
        stub_y1 = y + (5 * self.sin(angle + 90))
        stub_x2 = x - (5 * self.cos(angle + 90))
        stub_y2 = y - (5 * self.sin(angle + 90))

        if self.quadrant(stub_x1, stub_y1, stub_x2, stub_y2) >= 3:
            text_x = x + (10 * self.cos(angle + 90))
            text_y = y + (15 * self.sin(angle + 90))
        else:
            text_x = x - (10 * self.cos(angle + 90))
            text_y = y - (15 * self.sin(angle + 90))

        self.ctx.beginPath()
        self.ctx.moveTo(Pt(stub_x1), Pt(stub_y1))
        self.ctx.lineTo(Pt(stub_x2), Pt(stub_y2))
        self.ctx.strokeStyle = '#666'
        self.ctx.lineWidth = Pt(2)
        self.ctx.stroke()
        self.ctx.closePath()

        self.ctx.strokeText(text, Pt(text_x), Pt(text_y))

    def drawData(self, x, y):
        vertices = self.getVertices(x, y)

        for i in range(len(self.data)):
            next = i + 1 if i + 1 != len(self.data) else 0

            angle = atan2(vertices[i][1] - y, vertices[i][0] - x) * 180 / pi + 180;
            start_point = self.getPoint(x, y, vertices[i][0], vertices[i][1], angle, self.data[i] * 100 / ceil(max(self.data)))

            angle = atan2(vertices[next][1] - y, vertices[next][0] - x) * 180 / pi + 180;
            next_point = self.getPoint(x, y, vertices[next][0], vertices[next][1], angle, self.data[next] * 100 / ceil(max(self.data)))



            self.ctx.beginPath()
            self.ctx.moveTo(Pt(start_point['x']), Pt(start_point['y']))
            self.ctx.lineTo(Pt(next_point['x']), Pt(next_point['y']))

            self.ctx.strokeStyle = '#47a2f8'
            self.ctx.lineWidth = Pt(4)
            self.ctx.stroke()
            self.ctx.closePath()

            print("%s: %s" % (self.data[i], self.quadrant(x, y, vertices[i][0], vertices[i][1])))

            quadrant = self.quadrant(x, y, vertices[i][0], vertices[i][1])
            if quadrant is 1 or quadrant is 4:
                text_x = start_point['x'] + 15
                text_y = start_point['y'] - 15
            else:
                text_x = start_point['x'] - 30
                text_y = start_point['y'] - 15
            self.ctx.strokeText(self.data[i], Pt(text_x), Pt(text_y))

def draw_radar():

    print('Creating a blank presentation.')
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    print('Creating a Canvas context on the first slide.')
    canvas = Canvas(slide, prs.slide_height, prs.slide_width)
    ctx = canvas.getContext('2d')

    print('Creating Radar class instance.')
    radar = Radar(ctx, [ 2.7, 1.7, 2.5, 1.6, 0.6, 2.0, 1.8 ])

    print('Drawing Radar.')
    radar.draw(250, 250)

    print('Saving to `radar.pptx`')
    prs.save('radar.pptx')

if __name__ == '__main__':
    sys.path.append(path.abspath(path.join(path.dirname(__file__), '..')))
    from canvas import Canvas

    draw_radar()
