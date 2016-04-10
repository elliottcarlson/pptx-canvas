from pptx import Presentation
from pptx.util import Length, Pt
from os import sys, path

def draw_green_line():

    print('Creating a blank presentation.')
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    print('Creating a Canvas context on the first slide.')
    canvas = Canvas(slide, prs.slide_height, prs.slide_width)
    ctx = canvas.getContext('2d')

    print('Drawing a green line from 10pt,10pt to 100pt,100pt coordinates.')
    ctx.strokeStyle = '#00FF00'

    ctx.moveTo(Pt(10), Pt(10))
    ctx.lineTo(Pt(100), Pt(100))
    ctx.stroke()

    print('Saving to `green_line.pptx`')
    prs.save('green_line.pptx')

if __name__ == '__main__':
    sys.path.append(path.abspath(path.join(path.dirname(__file__), '..')))
    from canvas import Canvas

    draw_green_line()
