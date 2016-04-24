from pptx import Presentation
from pptx.util import Length, Pt
from os import sys, path
from math import atan2, ceil, cos, pi, pow, sin, sqrt
from numpy import arange
import operator
from collections import OrderedDict

DATA = OrderedDict([
    ('Social Media', 2.0),
    ('Direct Marketing', 1.8),
    ('Content', 1.4),
    ('Digital Advertising', 0.6),
    ('Multichannel', 0.0)
])

class BarChart(object):

    paddingTop = 10
    paddingBottom = 40
    paddingRight = 100
    paddingLeft = 30
    gridStep = 0.5

    def __init__(self, ctx, data):
        self.ctx = ctx
        self.data = data
        self.max_range = data[max(data.iteritems(), key=operator.itemgetter(1))[0]]

    def chartHeight(self, y):
        return y - self.paddingTop - self.paddingBottom

    def chartWidth(self, x):
        return x - self.paddingRight

    def gridHeight(self, y):
        return self.chartHeight(y) / (self.max_range / self.gridStep)

    def drawGrid(self, x, y):
        self.ctx.strokeStyle = '#e3e3e3'
        self.ctx.font = '8pt Georgia'

        j = self.paddingTop
        for i in arange(self.max_range, -0.5, -0.5):
            self.ctx.strokeText(i, Pt(5), Pt(j+3-10))
            self.ctx.beginPath()
            self.ctx.moveTo(Pt(self.paddingLeft), Pt(j))
            self.ctx.lineTo(Pt(self.chartWidth(x)), Pt(j))
            if i == 0:
                self.ctx.strokeStyle = '#333'
            self.ctx.stroke()
            self.ctx.closePath()

            self.ctx.save()
            self.ctx.font = 'italic 8pt Georgia'
            if (i).is_integer():
                self.ctx.strokeText('%d = average score' % i,
                    Pt(self.chartWidth(x) + 5), Pt(j+3-10))
            self.ctx.restore()
            self.ctx.closePath()
            j = j + self.gridHeight(y)
            self.ctx.font = '8pt Georgia'

    def barWidth(self, x):
        return (self.chartWidth(x) - self.paddingLeft) / len(self.data)

    def barCoords(self, i, val, x, y):
        chartHeight = self.paddingTop + self.chartHeight(y) -\
            (self.chartHeight(y) * val * 100 / self.max_range / 100)
        coords = {}
        coords['box_x1'] = self.paddingLeft + (self.barWidth(x) * i)
        coords['box_y1'] = self.paddingTop
        coords['box_x2'] = self.barWidth(x)
        coords['box_y2'] = self.chartHeight(y)
        coords['mid_x'] = self.paddingLeft + (self.barWidth(x) * i) +\
            (self.barWidth(x) / 2)
        coords['bar_x1'] = self.paddingLeft + (self.barWidth(x) * i) +\
            (self.barWidth(x) * 0.2)
        coords['bar_y1'] = chartHeight
        coords['bar_x2'] = self.barWidth(x) * 0.6
        coords['bar_y2'] = self.chartHeight(y) - chartHeight + self.paddingTop

        return coords

    def drawBars(self, x, y):
        i = 0
        for key, value in self.data.iteritems():
            coords = self.barCoords(i, value, x, y)
            self.ctx.fillStyle = '#999'
            self.ctx.fillRect(Pt(coords['bar_x1']), Pt(coords['bar_y1']),
                Pt(coords['bar_x2']), Pt(coords['bar_y2']))

            self.ctx.textAlign = 'center'
            self.ctx.fillStyle = '#fff'
            text_y = coords['bar_y1']
            text_y = text_y + 15

            if text_y > coords['box_y2']:
                text_y = coords['box_y2']
                self.ctx.fillStyle = '#999'

            self.ctx.strokeText(value, Pt(coords['mid_x'] - 10), Pt(text_y - 10))

            self.ctx.fillStyle = '#999'
            self.ctx.strokeText(key, Pt(coords['mid_x'] - 10), Pt(coords['box_y2'] +
                15))
#            self.wraptext(key, coords['mid_x'], coords['box_y2'] + 25,
#                    self.barWidth(x), 15)
            i = i + 1

    def draw(self, x, y):
        self.ctx.lineWidth = Pt(2)
        self.drawGrid(x, y)
        self.drawBars(x, y)


def draw_radar():

    print('Creating a blank presentation.')
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    print('Creating a Canvas context on the first slide.')
    canvas = Canvas(slide, prs.slide_height, prs.slide_width)
    ctx = canvas.getContext('2d')

    print('Creating BarChart class instance.')
    radar = BarChart(ctx, DATA)

    print('Drawing Bar Chart.')
    radar.draw(500, 225)

    print('Saving to `bar_chart.pptx`')
    prs.save('bar_chart.pptx')

if __name__ == '__main__':
    sys.path.append(path.abspath(path.join(path.dirname(__file__), '..')))
    from canvas import Canvas

    draw_radar()
