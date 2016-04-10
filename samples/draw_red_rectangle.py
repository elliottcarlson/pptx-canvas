from pptx import Presentation
from pptx.util import Length, Pt
from os import sys, path

def draw_red_rectangle():

    print('Creating a blank presentation.')
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    print('Creating a Canvas context on the first slide.')
    canvas = Canvas(slide, prs.slide_height, prs.slide_width)
    ctx = canvas.getContext('2d')

    print('Drawing a 150pt by 75pt red rectangle at 10pt,10pt coordinates.')
    ctx.fillStyle = '#FF0000'

    print ctx.fillStyle

    ctx.fillRect(Pt(10), Pt(10), Pt(150), Pt(75))

    print('Saving to `red_rectangle.pptx`')
    prs.save('red_rectangle.pptx')

if __name__ == '__main__':
    sys.path.append(path.abspath(path.join(path.dirname(__file__), '..')))
    from canvas import Canvas

    draw_red_rectangle()
