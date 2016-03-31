import os
import shutil
import tempfile
import unittest

from pptx import Presentation
from canvas.canvas import Canvas
from canvas.context.context import CanvasRenderingContext2D

class CanvasTest(unittest.TestCase):
    """Tests for canvas."""

    def setUp(self):
        # Create a pptx presentation with a single blank slide
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        self.slide_height = prs.slide_height
        self.slide_width = prs.slide_width

        self.canvas = Canvas(slide, self.slide_height, self.slide_width)

    def test_is_instance(self):
        self.assertIsInstance(self.canvas, Canvas)

    def test_has_correct_height(self):
        self.assertIs(self.slide_height, self.canvas._slide_height)

    def test_has_correct_width(self):
        self.assertIs(self.slide_width, self.canvas._slide_width)

    def test_gets_2d_context(self):
        ctx = self.canvas.getContext('2d')
        self.assertIsInstance(ctx, CanvasRenderingContext2D)

    def test_fails_unknown_context(self):
        ctx = self.canvas.getContext('webgl')
        self.assertIsNone(ctx)

    def test_supports_2d_context(self):
        supported = self.canvas.probablySupportsContext('2d')
        self.assertTrue(supported)

    def test_doesnt_supports_alt_context(self):
        supported = self.canvas.probablySupportsContext('webgl')
        self.assertFalse(supported)

if __name__ == "__main__" and __package__ is None:
    unittest.main()  # pragma: no cover
