import os
import shutil
import tempfile
import unittest

from pptx import Presentation
from canvas.canvas import Canvas
from canvas.context.context import CanvasRenderingContext2D

class CanvasStateTest(unittest.TestCase):
    """Tests for canvas."""

    def setUp(self):
        # Create a pptx presentation with a single blank slide
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        self.canvas = Canvas(slide, prs.slide_height, prs.slide_width)
        self.context = self.canvas.getContext('2d')

    def test_state(self):
        _shadowBlurStack = []
        self.context.shadowBlur = 1
        _shadowBlurStack.append(self.context.shadowBlur)
        self.context.save()
        self.context.shadowBlur = 10
        _shadowBlurStack.append(self.context.shadowBlur)
        self.context.save()
        self.context.shadowBlur = 12
        _shadowBlurStack.append(self.context.shadowBlur)
        self.context.save()
        self.context.shadowBlur = 15

        for v in reversed(_shadowBlurStack):
            self.context.restore()
            if self.context.shadowBlur is not v:
                self.fail('Incorrect value in restored stack')

        self.assertEqual(self.context.shadowBlur, 1)

if __name__ == "__main__" and __package__ is None:
    unittest.main()  # pragma: no cover
