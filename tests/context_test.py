import os
import shutil
import tempfile
import unittest

from pptx import Presentation
from canvas.canvas import Canvas
from canvas.context.context import CanvasRenderingContext2D

class ContextTest(unittest.TestCase):
    """Tests for canvas."""

    def setUp(self):
        # Create a pptx presentation with a single blank slide
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        self.canvas = Canvas(slide, prs.slide_height, prs.slide_width)
        self.context = self.canvas.getContext('2d')

    def test_is_instance(self):
        self.assertIsInstance(self.context, CanvasRenderingContext2D)

    def test_is_canvas_read_only_property(self):
        with self.assertRaises(AttributeError):
            self.context.canvas = None

if __name__ == "__main__" and __package__ is None:
    unittest.main()  # pragma: no cover
