import os
import shutil
import tempfile
import unittest

from pptx import Presentation
from canvas.canvas import Canvas
from canvas.context.context import CanvasRenderingContext2D

class CanvasTextDrawingStylesTest(unittest.TestCase):
    def setUp(self):
        # Create a pptx presentation with a single blank slide
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        self.canvas = Canvas(slide, prs.slide_height, prs.slide_width)
        self.context = self.canvas.getContext('2d')

    def test_textAlign(self):
        # Default value
        self.assertEqual(self.content.textAlign, 'start')

        # Change to valid value
        self.context.textAlign = 'center'
        self.assertEqual(self.content.textAlign, 'center')

        # Attempt to change to invalid value
        self.context.textAlign = 'invalid'
        self.assertNotEqual(self.content.textAlign, 'invalid')

    def test_textBaseline(self):
        # Default value
        self.assertEqual(self.content.direction, 'alphabetic')

        # Change to valid value
        self.context.textBaseline = 'hanging'
        self.assertEqual(self.content.textBaseline, 'hanging')

        # Attempt to change to invalid value
        self.context.textBaseline = 'invalid'
        self.assertNotEqual(self.content.textBaseline, 'invalid')

    def test_direction(self):
        # Default value
        self.assertEqual(self.content.direction, 'inherit')

        # Change to valid value
        self.context.direction = 'ltr'
        self.assertEqual(self.content.direction, 'ltr')

        # Attempt to change to invalid value
        self.context.direction = 'invalid'
        self.assertNotEqual(self.content.direction, 'invalid')

if __name__ == "__main__" and __package__ is None:
    unittest.main()  # pragma: no cover
